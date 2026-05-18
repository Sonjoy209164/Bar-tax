from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "engineering_architecture_slide.pptx"
DOWNLOAD = ROOT / "frontend" / "downloads" / "engineering_architecture_slide.pptx"


NAVY = RGBColor(23, 32, 51)
MUTED = RGBColor(86, 99, 119)
LIGHT = RGBColor(247, 249, 252)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(204, 214, 228)
TEAL = RGBColor(29, 138, 165)
GREEN = RGBColor(47, 143, 91)
ORANGE = RGBColor(213, 123, 31)
PURPLE = RGBColor(98, 102, 189)
MAGENTA = RGBColor(176, 99, 155)
GRAY = RGBColor(105, 115, 134)
RED = RGBColor(190, 75, 75)


def textbox(slide, x, y, w, h, text, *, size=12, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def card(slide, x, y, w, h, title, body, *, accent=TEAL, body_size=8.6):
    shadow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.025), y + Inches(0.035), w, h)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(218, 225, 236)
    shadow.line.fill.background()
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = LINE
    s.line.width = Pt(0.8)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.055))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, x + Inches(0.1), y + Inches(0.09), w - Inches(0.2), Inches(0.18), title, size=9.8, bold=True)
    textbox(slide, x + Inches(0.1), y + Inches(0.31), w - Inches(0.2), h - Inches(0.36), body, size=body_size, color=MUTED)
    return s


def label(slide, x, y, w, h, text, *, color=TEAL, size=11):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    tf = s.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(255, 255, 255)
    return s


def arrow(slide, x1, y1, x2, y2, *, color=LINE):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True
    return line


def flow_node(slide, x, y, w, h, title, detail, *, color):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = color
    s.line.width = Pt(1.3)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.06))
    top.fill.solid()
    top.fill.fore_color.rgb = color
    top.line.fill.background()
    textbox(slide, x + Inches(0.07), y + Inches(0.08), w - Inches(0.14), Inches(0.18), title, size=8.8, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, x + Inches(0.08), y + Inches(0.29), w - Inches(0.16), h - Inches(0.32), detail, size=7.2, color=MUTED, align=PP_ALIGN.CENTER)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()

    textbox(slide, Inches(0.35), Inches(0.18), Inches(9.2), Inches(0.34), "Engineering Architecture: Inventory RAG Chatbot", size=22, bold=True)
    textbox(slide, Inches(0.37), Inches(0.58), Inches(12.3), Inches(0.28), "Board-ready view of runtime flow, code ownership, technology choices, and prompt boundaries.", size=10.5, color=MUTED)

    # Runtime flow.
    label(slide, Inches(0.38), Inches(0.98), Inches(2.35), Inches(0.32), "Runtime Request Flow", color=NAVY, size=10.5)
    x0 = Inches(0.38)
    y0 = Inches(1.46)
    w = Inches(1.22)
    h = Inches(0.76)
    nodes = [
        ("UI", "chat.html\nchat.js", GREEN),
        ("API", "FastAPI\nroutes", TEAL),
        ("Service", "InventoryService\norchestrator", NAVY),
        ("Understand", "intent\nslots\nmemory", PURPLE),
        ("Retrieve", "structured\nvector\nlexical", ORANGE),
        ("Decide", "ranker\nplanner\nevidence", MAGENTA),
        ("Write", "template or\nLLM writer", GRAY),
        ("Verify", "claim + fit\nchecks", RED),
    ]
    for i, (title, detail, color) in enumerate(nodes):
        x = x0 + i * Inches(1.55)
        flow_node(slide, x, y0, w, h, title, detail, color=color)
        if i < len(nodes) - 1:
            arrow(slide, x + w + Inches(0.04), y0 + h / 2, x + Inches(1.48), y0 + h / 2)

    # Ownership matrix.
    label(slide, Inches(0.38), Inches(2.52), Inches(2.15), Inches(0.32), "Layer Ownership", color=NAVY, size=10.5)
    col_w = Inches(2.02)
    row_h = Inches(1.18)
    left = Inches(0.38)
    top = Inches(3.0)
    cards = [
        ("Data", "catalog.jsonl\nstorage.py\nschemas.py\n\nOwns truth:\nprice, stock, SKU,\nattributes", TEAL),
        ("Sync/API", "pos_sync.py\nroutes_inventory.py\nmain.py\n\nOwns validation,\nserving, auth", GREEN),
        ("Understanding", "fashion_retail.py\nintent.py\nllm_* files\n\nOwns language,\nintent, slots", PURPLE),
        ("Retrieval", "embedder.py\nvector_store_base.py\nelasticsearch_store.py\nbm25_index.py", ORANGE),
        ("Decision", "reranker.py\ndecisioning.py\nplanner.py\nevidence_contract.py", MAGENTA),
        ("Answer", "natural_answer.py\nverifier.py\nanswer_critic.py\nchat.js", GRAY),
    ]
    for i, (title, body, color) in enumerate(cards):
        card(slide, left + i * Inches(2.12), top, col_w, row_h, title, body, accent=color, body_size=7.6)

    # Prompt boundaries.
    label(slide, Inches(0.38), Inches(4.55), Inches(2.3), Inches(0.32), "Prompt Boundaries", color=NAVY, size=10.5)
    card(slide, Inches(0.38), Inches(5.02), Inches(3.95), Inches(1.35), "Prompts Can Do", "Classify intent/slots as JSON\nPlan from recent conversation\nPick from bounded candidates\nRewrite supported facts warmly\nCritique answer against facts", accent=PURPLE, body_size=8.4)
    card(slide, Inches(4.58), Inches(5.02), Inches(3.95), Inches(1.35), "Prompts Cannot Do", "Invent products, stock, price, policy\nOverride evidence contract\nRecommend excluded products\nTreat cross-sells as substitutes\nIgnore verifier failure", accent=RED, body_size=8.4)
    card(slide, Inches(8.78), Inches(5.02), Inches(3.95), Inches(1.35), "Key Prompt Files", "llm_intent_classifier.py\nllm_slot_extractor.py\nintent_planner.py\nllm_reasoner.py\nnatural_answer.py\nanswer_critic.py", accent=TEAL, body_size=8.0)

    # Bottom note.
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(6.58), Inches(12.35), Inches(0.42))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(236, 243, 248)
    box.line.color.rgb = LINE
    textbox(
        slide,
        Inches(0.55),
        Inches(6.69),
        Inches(12.0),
        Inches(0.18),
        "Engineering thesis: deterministic layers own correctness; retrieval owns evidence; LLM layers improve interpretation and wording under strict contracts.",
        size=9.2,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )

    return prs


def main() -> None:
    prs = build()
    prs.save(OUT)
    DOWNLOAD.parent.mkdir(parents=True, exist_ok=True)
    prs.save(DOWNLOAD)
    print(OUT)
    print(DOWNLOAD)


if __name__ == "__main__":
    main()
