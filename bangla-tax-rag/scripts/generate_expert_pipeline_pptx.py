from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from PIL import Image


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "expert_pipeline_presentation.pptx"
FULL_PIPELINE = ROOT / "docs" / "assets" / "full_current_inventory_pipeline.png"
TREE_PIPELINE = ROOT / "docs" / "assets" / "inventory_pipeline_tree.png"


NAVY = RGBColor(23, 32, 51)
MUTED = RGBColor(88, 101, 121)
TEAL = RGBColor(29, 138, 165)
GREEN = RGBColor(47, 143, 91)
ORANGE = RGBColor(213, 123, 31)
PURPLE = RGBColor(98, 102, 189)
MAGENTA = RGBColor(176, 99, 155)
GRAY = RGBColor(105, 115, 134)
LIGHT = RGBColor(247, 249, 252)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(204, 214, 228)


def add_bg(slide) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(12.35), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.47), Inches(0.82), Inches(12.1), Inches(0.34))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12.5)
        sp.font.color.rgb = MUTED


def add_footer(slide, idx: int) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(7.13), Inches(12.45), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.48), Inches(7.17), Inches(12.0), Inches(0.22))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Inventory RAG Pipeline | bangla-tax-rag | slide {idx}"
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED


def add_bullets(slide, x, y, w, h, bullets: list[str], *, size: float = 17, color=NAVY, gap: float = 1.0) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(5 * gap)


def add_card(slide, x, y, w, h, title: str, bullets: list[str], *, accent=TEAL, title_size=18, body_size=12.5) -> None:
    shadow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.04), y + Inches(0.05), w, h)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(218, 225, 236)
    shadow.line.fill.background()
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    title_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.16), w - Inches(0.35), Inches(0.32))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(title_size)
    p.font.color.rgb = NAVY
    body = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.58), w - Inches(0.38), h - Inches(0.72))
    btf = body.text_frame
    btf.clear()
    btf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(body_size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(4)


def add_quote(slide, text: str, x, y, w, h, *, accent=TEAL) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(236, 243, 248)
    shape.line.color.rgb = LINE
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.18), w - Inches(0.35), h - Inches(0.26))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(15)
    p.font.color.rgb = NAVY


def add_image_fit(slide, image_path: Path, x, y, w, h) -> None:
    with Image.open(image_path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pw = int(iw * scale)
    ph = int(ih * scale)
    px = x + int((w - pw) / 2)
    py = y + int((h - ph) / 2)
    slide.shapes.add_picture(str(image_path), px, py, width=pw, height=ph)


def add_flow(slide, labels: list[tuple[str, RGBColor]], x, y, w, h) -> None:
    gap = Inches(0.12)
    item_w = (w - gap * (len(labels) - 1)) / len(labels)
    for i, (label, color) in enumerate(labels):
        bx = x + i * (item_w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, bx, y, item_w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)


def add_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle)
    return slide


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = add_slide(prs, "Current Inventory RAG Pipeline", "Catalog-grounded boutique retail chatbot for customer query answering")
    add_quote(
        slide,
        "Core idea: the LLM does not own truth. Catalog data, retrieval, decision rules, evidence contracts, and verification control what the bot can safely say.",
        Inches(0.75), Inches(1.55), Inches(11.8), Inches(1.05), accent=TEAL,
    )
    add_flow(
        slide,
        [
            ("Catalog", TEAL),
            ("Understand", PURPLE),
            ("Retrieve", ORANGE),
            ("Decide", MAGENTA),
            ("Write", GRAY),
            ("Verify", NAVY),
        ],
        Inches(0.75), Inches(3.05), Inches(11.8), Inches(0.72),
    )
    add_bullets(
        slide,
        Inches(1.0), Inches(4.25), Inches(11.4), Inches(1.6),
        [
            "Built for sarees, bags, cosmetics, jewelry, watches, shoes, panjabi, shirts, pants, perfumes, and other boutique catalog items.",
            "Handles Bangla, Banglish, English, and mixed customer wording.",
            "Answers exact inventory questions without inventing product facts.",
        ],
        size=17,
    )
    add_footer(slide, 1)

    # 2
    slide = add_slide(prs, "What The Bot Can Do Now", "Customer-facing capabilities in the current architecture")
    add_card(slide, Inches(0.65), Inches(1.35), Inches(3.8), Inches(2.05), "Catalog Questions", [
        "Availability: do you have this item?",
        "Price, stock, SKU, product details",
        "Same design in another color",
        "Size availability",
    ], accent=TEAL)
    add_card(slide, Inches(4.75), Inches(1.35), Inches(3.8), Inches(2.05), "Shopping Help", [
        "Budget-based product search",
        "Occasion-based suggestions",
        "Styling advice and accessory match",
        "Compare products or fabrics",
    ], accent=PURPLE)
    add_card(slide, Inches(8.85), Inches(1.35), Inches(3.8), Inches(2.05), "Operations", [
        "Delivery/payment/refund/exchange policy QA",
        "Order/cart workflow hooks",
        "Image-search pathway",
        "Feedback and trace capture",
    ], accent=GREEN)
    add_quote(slide, "Best current use case: answering boutique customer questions from current inventory with grounded, human-friendly replies.", Inches(0.65), Inches(4.2), Inches(12.0), Inches(0.9), accent=ORANGE)
    add_footer(slide, 2)

    # 3
    slide = add_slide(prs, "Architecture Positioning", "This is not a prompt-only chatbot")
    add_card(slide, Inches(0.65), Inches(1.3), Inches(3.9), Inches(4.9), "What Owns Truth", [
        "Catalog JSONL / SQLite mirror",
        "Pydantic product schemas",
        "Policy JSON",
        "Order and sync records",
        "Evidence contract",
    ], accent=TEAL, body_size=14)
    add_card(slide, Inches(4.75), Inches(1.3), Inches(3.9), Inches(4.9), "What Owns Search", [
        "Structured filters",
        "Bangla/Banglish normalization",
        "Vector retrieval",
        "Lexical/SKU search",
        "Reranking and scoring",
    ], accent=ORANGE, body_size=14)
    add_card(slide, Inches(8.85), Inches(1.3), Inches(3.9), Inches(4.9), "What The LLM Owns", [
        "Optional intent/slot JSON extraction",
        "Optional conversation planning",
        "Optional candidate reasoning",
        "Natural answer writing",
        "Optional answer critique",
    ], accent=PURPLE, body_size=14)
    add_footer(slide, 3)

    # 4
    slide = add_slide(prs, "Full Pipeline Diagram", "End-to-end flow from catalog setup to verified customer answer")
    add_image_fit(slide, FULL_PIPELINE, Inches(0.45), Inches(1.18), Inches(12.45), Inches(5.65))
    add_footer(slide, 4)

    # 5
    slide = add_slide(prs, "Tree View", "Responsibility boundaries by subsystem")
    add_image_fit(slide, TREE_PIPELINE, Inches(0.35), Inches(1.1), Inches(12.65), Inches(5.75))
    add_footer(slide, 5)

    # 6
    slide = add_slide(prs, "Layer 1: Data And Storage", "How catalogs are saved and trusted")
    add_card(slide, Inches(0.65), Inches(1.25), Inches(3.8), Inches(4.95), "Catalog DB", [
        "data/inventory/catalog.jsonl",
        "InventoryItemRecord in app/core/schemas.py",
        "Fields: product_id, SKU, category, price, stock, attributes, tags, status",
        "This is the source of truth for product facts",
    ], accent=TEAL)
    add_card(slide, Inches(4.75), Inches(1.25), Inches(3.8), Inches(4.95), "Mirror Store", [
        "app/inventory/storage.py",
        "JSONL default; SQLite mirror available",
        "Pydantic validation before persistence",
        "Business signals stored separately from catalog facts",
    ], accent=GREEN)
    add_card(slide, Inches(8.85), Inches(1.25), Inches(3.8), Inches(4.95), "Sync", [
        "app/inventory/pos_sync.py",
        "CSV and webhook import",
        "Merge by product id/SKU",
        "Audit logs and rebuild status track drift",
    ], accent=ORANGE)
    add_footer(slide, 6)

    # 7
    slide = add_slide(prs, "Layer 2: Query Understanding", "How the bot reads Bangla, Banglish, English, and mixed text")
    add_card(slide, Inches(0.65), Inches(1.25), Inches(3.8), Inches(4.95), "Normalization", [
        "banglish_normalizer.py",
        "fuzzy_corrector.py",
        "Bangla digit translation",
        "Alias and typo expansion before slot extraction",
    ], accent=TEAL)
    add_card(slide, Inches(4.75), Inches(1.25), Inches(3.8), Inches(4.95), "Intent", [
        "intent.py deterministic classifier",
        "llm_intent_classifier.py optional JSON classifier",
        "intent_planner.py optional multi-turn planner",
        "Outputs confidence and ambiguity reason",
    ], accent=PURPLE)
    add_card(slide, Inches(8.85), Inches(1.25), Inches(3.8), Inches(4.95), "Slots", [
        "fashion_retail.py regex and aliases",
        "llm_slot_extractor.py optional JSON extraction",
        "Finds category, color, size, budget, fabric, occasion, gender, design id",
        "Regex wins when safer",
    ], accent=GREEN)
    add_footer(slide, 7)

    # 8
    slide = add_slide(prs, "Layer 3: Retrieval", "Three retrieval paths, each solving a different failure mode")
    add_card(slide, Inches(0.65), Inches(1.25), Inches(3.8), Inches(4.95), "Structured Search", [
        "Best for exact catalog facts",
        "Filters category, color, size, budget, design id, stock",
        "Handles same design and size availability",
        "No prompt involved",
    ], accent=TEAL)
    add_card(slide, Inches(4.75), Inches(1.25), Inches(3.8), Inches(4.95), "Vector Search", [
        "embedder.py builds query/product vectors",
        "local_store.py for local JSONL vector storage",
        "elasticsearch_store.py for kNN + filters",
        "Pinecone and Milvus adapters available",
    ], accent=ORANGE)
    add_card(slide, Inches(8.85), Inches(1.25), Inches(3.8), Inches(4.95), "Lexical Search", [
        "BM25/lexical matching",
        "SKU, product id, brand, category, exact name",
        "Elasticsearch multi_match path",
        "Protects against embedding fuzziness",
    ], accent=GREEN)
    add_footer(slide, 8)

    # 9
    slide = add_slide(prs, "Layer 4: Decisions And Evidence", "Where recommendations become controlled decisions")
    add_card(slide, Inches(0.65), Inches(1.25), Inches(3.8), Inches(4.95), "Ranking", [
        "reranker.py computes ecommerce fit",
        "Signals: semantic, lexical, exact SKU/name, category, brand, product type",
        "Price fit, stock fit, metadata/spec fit",
        "Penalizes unrelated and out-of-stock candidates",
    ], accent=ORANGE)
    add_card(slide, Inches(4.75), Inches(1.25), Inches(3.8), Inches(4.95), "Answer Plan", [
        "planner.py enriches the recommendation plan",
        "Chooses primary, alternatives, cross-sells, next question",
        "Adds risk notes and tradeoffs",
        "Defines what the writer may say",
    ], accent=MAGENTA)
    add_card(slide, Inches(8.85), Inches(1.25), Inches(3.8), Inches(4.95), "Evidence Contract", [
        "evidence_contract.py is the safety boundary",
        "Allowed claims, missing facts, contradictions",
        "Rejected candidate ids",
        "Prevents fluent hallucination",
    ], accent=PURPLE)
    add_footer(slide, 9)

    # 10
    slide = add_slide(prs, "Layer 5: Answer Writing And Verification", "How a safe answer becomes human-friendly")
    add_card(slide, Inches(0.65), Inches(1.25), Inches(3.8), Inches(4.95), "Deterministic Writer", [
        "fashion_retail.py templates",
        "Safe for exact stock, price, size, color, design id",
        "Works when LLM is unavailable",
        "Bangla/Banglish localization support",
    ], accent=TEAL)
    add_card(slide, Inches(4.75), Inches(1.25), Inches(3.8), Inches(4.95), "Natural Writer", [
        "natural_answer.py",
        "inventory_service._build_inventory_answer_messages()",
        "Ollama/OpenAI-compatible writer",
        "Must obey answer_plan and writer_contract",
    ], accent=PURPLE)
    add_card(slide, Inches(8.85), Inches(1.25), Inches(3.8), Inches(4.95), "Verifier", [
        "verifier.py checks final answer",
        "answer_critic.py optional LLM critic",
        "Rejects unsupported products, prices, stock claims",
        "Allows one cautious retry",
    ], accent=GRAY)
    add_footer(slide, 10)

    # 11
    slide = add_slide(prs, "Prompt Registry", "Where prompts exist and what each prompt is allowed to decide")
    add_card(slide, Inches(0.55), Inches(1.18), Inches(4.0), Inches(2.2), "Understanding Prompts", [
        "llm_intent_classifier.py: JSON intent, slots, confidence",
        "llm_slot_extractor.py: JSON shopping slots",
        "intent_planner.py: multi-turn plan and constraints",
    ], accent=PURPLE, body_size=11.5)
    add_card(slide, Inches(4.75), Inches(1.18), Inches(4.0), Inches(2.2), "Decision Prompts", [
        "llm_reasoner.py: pick product ids only from candidates",
        "Never invent products",
        "Can return none_fit",
    ], accent=MAGENTA, body_size=11.5)
    add_card(slide, Inches(8.95), Inches(1.18), Inches(3.75), Inches(2.2), "Writing Prompts", [
        "natural_answer.py: warm boutique answer",
        "inventory_service.py: strict ecommerce writer contract",
        "answer_critic.py: quality review",
    ], accent=GREEN, body_size=11.5)
    add_quote(slide, "Prompt governance rule: prompts may interpret or phrase, but they cannot override catalog facts, selected product roles, excluded product ids, or verification.", Inches(0.7), Inches(4.3), Inches(11.9), Inches(0.9), accent=ORANGE)
    add_footer(slide, 11)

    # 12
    slide = add_slide(prs, "Technology Stack", "What technology is used at each stage")
    add_card(slide, Inches(0.65), Inches(1.25), Inches(3.8), Inches(4.95), "Backend", [
        "FastAPI",
        "Pydantic",
        "PyYAML settings",
        "httpx",
        "pytest",
    ], accent=GREEN, body_size=16)
    add_card(slide, Inches(4.75), Inches(1.25), Inches(3.8), Inches(4.95), "Retrieval", [
        "JSONL / SQLite mirror",
        "rank-bm25",
        "Transformers / multilingual / OpenAI embeddings",
        "Local vector store",
        "Elasticsearch / Pinecone / Milvus adapters",
    ], accent=ORANGE, body_size=14)
    add_card(slide, Inches(8.85), Inches(1.25), Inches(3.8), Inches(4.95), "Generation", [
        "Ollama qwen3:8b prompts",
        "OpenAI-compatible chat client",
        "Deterministic fallback templates",
        "Evidence contract",
        "Final verifier",
    ], accent=PURPLE, body_size=14)
    add_footer(slide, 12)

    # 13
    slide = add_slide(prs, "Runtime Endpoints", "How to demo and inspect the system")
    add_card(slide, Inches(0.65), Inches(1.25), Inches(3.8), Inches(4.95), "User UI", [
        "http://127.0.0.1:4849/frontend/chat.html",
        "Chat with the bot",
        "Open/close catalog panel",
        "Send feedback and order actions",
    ], accent=GREEN)
    add_card(slide, Inches(4.75), Inches(1.25), Inches(3.8), Inches(4.95), "API Docs", [
        "http://127.0.0.1:4849/docs",
        "Manual endpoint testing",
        "Inventory ask/search/sync routes",
        "Order, feedback, image routes",
    ], accent=TEAL)
    add_card(slide, Inches(8.85), Inches(1.25), Inches(3.8), Inches(4.95), "Core Calls", [
        "GET /inventory/items",
        "POST /inventory/ask",
        "POST /inventory/search",
        "POST /inventory/sync/rebuild",
        "GET /inventory/chat-traces/{trace_id}",
    ], accent=ORANGE, body_size=12.5)
    add_footer(slide, 13)

    # 14
    slide = add_slide(prs, "Safeguards And Failure Handling", "What prevents bad customer answers")
    add_bullets(
        slide,
        Inches(0.85), Inches(1.35), Inches(11.8), Inches(1.1),
        [
            "If the catalog has no matching product, the bot should abstain or ask a focused clarification.",
            "If the answer writer tries to add unsupported stock, price, policy, product, or discount claims, verifier catches it.",
            "If Ollama or an LLM prompt fails, the deterministic pipeline still produces a safer answer.",
        ],
        size=18,
    )
    add_card(slide, Inches(0.85), Inches(3.15), Inches(3.7), Inches(2.25), "Risk", ["Vague customer query", "Weak retrieval", "LLM overclaims"], accent=ORANGE, body_size=14)
    add_card(slide, Inches(4.85), Inches(3.15), Inches(3.7), Inches(2.25), "Control", ["Clarification gate", "Evidence contract", "Verifier"], accent=PURPLE, body_size=14)
    add_card(slide, Inches(8.85), Inches(3.15), Inches(3.7), Inches(2.25), "Outcome", ["Ask a better question", "Return supported facts", "Abstain if unsafe"], accent=GREEN, body_size=14)
    add_footer(slide, 14)

    # 15
    slide = add_slide(prs, "Honest Limitations", "Important to say clearly to an expert engineer")
    add_bullets(
        slide,
        Inches(0.75), Inches(1.25), Inches(12.0), Inches(4.7),
        [
            "Image matching exists as a pathway, but production-grade visual similarity needs a stronger image embedding/index pipeline.",
            "Customer memory exists, but profile memory should be privacy-scoped and consent-aware before serious deployment.",
            "Real-time POS sync is currently import/webhook style; true POS integration needs source-specific connectors and conflict policy.",
            "Large multi-brand catalogs require stronger taxonomy governance, synonym management, and monitoring for stale vectors.",
            "The answer quality ceiling depends heavily on catalog completeness: missing sizes, colors, variants, or policies cannot be invented safely.",
        ],
        size=18,
    )
    add_footer(slide, 15)

    # 16
    slide = add_slide(prs, "Recommended Next Build Priorities", "A pragmatic roadmap toward a production-grade human-like assistant")
    add_card(slide, Inches(0.65), Inches(1.25), Inches(3.8), Inches(4.95), "1. Data Quality", [
        "Normalize variant groups and design ids",
        "Require size/color/fabric fields by category",
        "Add policy coverage tests",
        "Add catalog quality dashboard",
    ], accent=TEAL)
    add_card(slide, Inches(4.75), Inches(1.25), Inches(3.8), Inches(4.95), "2. Retrieval Quality", [
        "Use real multilingual embeddings",
        "Benchmark Elasticsearch vs local vector store",
        "Add synonym/taxonomy management",
        "Expand multilingual eval set",
    ], accent=ORANGE)
    add_card(slide, Inches(8.85), Inches(1.25), Inches(3.8), Inches(4.95), "3. Conversation Quality", [
        "Improve customer profile memory",
        "Use stronger answer writer evals",
        "Add image-based matching tests",
        "Run build-test-fix loop weekly",
    ], accent=PURPLE)
    add_footer(slide, 16)

    return prs


def main() -> None:
    missing = [str(p) for p in (FULL_PIPELINE, TREE_PIPELINE) if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing diagram image(s): " + ", ".join(missing))
    prs = build_deck()
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
