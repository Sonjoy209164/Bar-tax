from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TREE_IMAGE = ROOT / "docs" / "assets" / "inventory_pipeline_tree.png"
OUTPUT = ROOT / "docs" / "inventory_pipeline_tree_7_slide_deck.pptx"


INK = RGBColor(24, 31, 45)
MUTED = RGBColor(91, 103, 121)
LIGHT = RGBColor(241, 245, 249)
LINE = RGBColor(203, 213, 225)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)
PURPLE = RGBColor(124, 58, 237)


def add_textbox(slide, x, y, w, h, text="", *, font_size=18, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None, accent=BLUE) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_textbox(slide, 0.55, 0.28, 10.4, 0.48, title, font_size=25, bold=True)
    if subtitle:
        add_textbox(slide, 0.58, 0.78, 11.8, 0.42, subtitle, font_size=10.5, color=MUTED)
    line = slide.shapes.add_shape(1, Inches(0.58), Inches(1.18), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.color.rgb = accent


def add_section_label(slide, x, y, text: str, accent=BLUE) -> None:
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(2.15), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.color.rgb = accent
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(10.5)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_bullets(slide, x, y, w, h, bullets: list[str], *, font_size=14.5, color=INK) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(7)
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p._p.get_or_add_pPr().set("marL", "220000")
        p._p.get_or_add_pPr().set("indent", "-140000")


def add_card(slide, x, y, w, h, heading: str, bullets: list[str], *, accent=BLUE) -> None:
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = LINE
    card.line.width = Pt(1.2)
    add_textbox(slide, x + 0.18, y + 0.14, w - 0.36, 0.28, heading, font_size=13.5, bold=True, color=accent)
    add_bullets(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.56, bullets, font_size=10.5)


def add_tree_thumb(slide, x=8.45, y=1.45, w=4.15) -> None:
    slide.shapes.add_picture(str(TREE_IMAGE), Inches(x), Inches(y), width=Inches(w))
    add_textbox(slide, x, y + 2.73, w, 0.28, "Pipeline tree reference", font_size=8.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_footer(slide, slide_no: int) -> None:
    add_textbox(
        slide,
        0.58,
        7.12,
        7.2,
        0.22,
        "Inventory RAG chatbot | catalog-grounded retail assistant",
        font_size=8.5,
        color=MUTED,
    )
    add_textbox(slide, 12.1, 7.12, 0.7, 0.22, str(slide_no), font_size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = [
        {
            "title": "Inventory RAG Pipeline: Full Tree View",
            "subtitle": "The bot is not one prompt. It is a layered system: catalog facts -> retrieval -> decisions -> grounded answer.",
            "accent": BLUE,
            "type": "cover",
        },
        {
            "title": "Layer 1: DB Setup + Catalog Ground Truth",
            "subtitle": "Where product truth lives before the model ever sees a question.",
            "accent": GREEN,
            "cards": [
                (
                    "What It Does",
                    [
                        "Stores products as structured facts: product_id, SKU, name, category, price, stock, status, attributes, tags.",
                        "Keeps catalog, order, feedback, policy, and sync audit data separate so each source has a clear responsibility.",
                        "Treats catalog/POS facts as the source of truth; the LLM cannot invent product availability.",
                    ],
                ),
                (
                    "Technology Used",
                    [
                        "JSONL catalog: data/inventory/catalog.jsonl.",
                        "Optional SQLite mirror: data/inventory/inventory_mirror.sqlite3.",
                        "Pydantic schemas validate product/order/API objects.",
                    ],
                ),
                (
                    "Logic",
                    [
                        "Every item must carry searchable attributes like color, size, design_id, gender, fabric, occasion, and stock.",
                        "include_in_rag/status decide whether a product should appear in answers.",
                        "Hard fields beat natural language: exact stock, price, SKU, and size are never guessed.",
                    ],
                ),
            ],
        },
        {
            "title": "Layer 2: Ingest, Sync + Search Document Build",
            "subtitle": "How raw shop data becomes searchable evidence.",
            "accent": ORANGE,
            "cards": [
                (
                    "What It Does",
                    [
                        "Imports or updates products from manual JSONL, POS CSV-style imports, or webhook payloads.",
                        "Normalizes fields and writes audit records so stock/catalog changes can be traced.",
                        "Turns each product into searchable text plus exact metadata for filters.",
                    ],
                ),
                (
                    "Technology Used",
                    [
                        "FastAPI sync endpoints under inventory routes.",
                        "Pydantic validation and Python ingestion logic.",
                        "Audit storage: data/inventory/sync_audit.jsonl.",
                    ],
                ),
                (
                    "Logic",
                    [
                        "Merge/update by product_id or SKU to avoid duplicate inventory records.",
                        "Build search text from name, category, brand, attributes, tags, and description.",
                        "Use deterministic IDs so rebuilds replace old vector records cleanly.",
                    ],
                ),
            ],
        },
        {
            "title": "Layer 3: Query Understanding",
            "subtitle": "How Bangla, Banglish, and English questions become structured shopping intent.",
            "accent": PURPLE,
            "cards": [
                (
                    "What It Does",
                    [
                        "Detects whether the user wants product search, size check, same-design color, styling, policy, order, or small talk.",
                        "Extracts slots: category, color, fabric, size, brand, budget, occasion, gender, and in-stock intent.",
                        "Reads recent conversation when the customer says follow-up phrases like 'same design blue ache?'.",
                    ],
                ),
                (
                    "Technology Used",
                    [
                        "Regex and alias dictionaries for deterministic retail terms.",
                        "Banglish normalizer and fuzzy corrector for typo-heavy customer text.",
                        "Optional Ollama qwen3:8b prompts for intent, slot extraction, and conversation planning.",
                    ],
                ),
                (
                    "Logic",
                    [
                        "Deterministic extraction wins for hard fields; LLM output fills gaps only when useful.",
                        "Low-confidence or vague questions trigger one clarifying question instead of random recommendations.",
                        "Memory/context helps resolve references to last shown products and prior preferences.",
                    ],
                ),
            ],
        },
        {
            "title": "Layer 4: Retrieval + Ranking",
            "subtitle": "How the system finds the right products from the catalog.",
            "accent": BLUE,
            "cards": [
                (
                    "What It Does",
                    [
                        "Runs structured filters for exact requirements: category, size, color, design_id, price, gender, stock.",
                        "Runs semantic/vector search for flexible questions like 'elegant eid saree under 5000'.",
                        "Combines lexical, semantic, and retail-fit signals before choosing candidates.",
                    ],
                ),
                (
                    "Technology Used",
                    [
                        "Embedding layer: transformers/multilingual/deterministic/OpenAI-compatible modes.",
                        "Vector store: local JSONL by default; Elasticsearch adapter exists for dense vector + filter search.",
                        "BM25/lexical search catches exact product names, SKUs, brands, and categories.",
                    ],
                ),
                (
                    "Logic",
                    [
                        "Structured filters protect hard correctness; embeddings help recall but do not override facts.",
                        "Reranking rewards exact category/slot match, stock fit, budget fit, and product-name match.",
                        "Optional LLM reasoner can select only from bounded retrieved candidates, not the whole world.",
                    ],
                ),
            ],
        },
        {
            "title": "Layer 5: Evidence Contract + Decision Plan",
            "subtitle": "The anti-hallucination layer between retrieval and final answer.",
            "accent": GREEN,
            "cards": [
                (
                    "What It Does",
                    [
                        "Turns retrieved products into an evidence package: allowed facts, missing facts, risks, tradeoffs, and product roles.",
                        "Builds an answer plan: primary product, alternatives, cross-sells, exclusions, caveats, next question.",
                        "Decides answer vs clarify vs abstain when evidence is weak or unsupported.",
                    ],
                ),
                (
                    "Technology Used",
                    [
                        "Pydantic evidence contract models.",
                        "Deterministic planner, verifier, and policy gates.",
                        "Optional answer_plan enrichment before natural language writing.",
                    ],
                ),
                (
                    "Logic",
                    [
                        "Primary and alternative products must come from retrieved catalog evidence.",
                        "Cross-sell products are add-ons only, never replacements.",
                        "Excluded or out-of-stock products cannot be pitched as available.",
                    ],
                ),
            ],
        },
        {
            "title": "Layer 6: Answer Writing, UI + Feedback Loop",
            "subtitle": "How the final response reaches the customer and improves over time.",
            "accent": ORANGE,
            "cards": [
                (
                    "What It Does",
                    [
                        "Writes a human-friendly answer in the customer's language style: Bangla, Banglish, or English.",
                        "Serves the chat UI, hideable catalog panel, product actions, order flow, and API docs.",
                        "Stores thumbs up/down, comments, trace IDs, failed cases, and response evidence for tuning.",
                    ],
                ),
                (
                    "Technology Used",
                    [
                        "FastAPI backend: /frontend/chat.html and /docs.",
                        "Plain HTML/CSS/JavaScript frontend.",
                        "Optional Ollama/OpenAI-compatible natural writer plus deterministic fallback templates.",
                    ],
                ),
                (
                    "Logic",
                    [
                        "Strict writer prompt follows answer_plan and writer_contract; it returns structured JSON.",
                        "Answer critic/verifier checks product names, prices, stock, exclusions, and unsupported claims.",
                        "Feedback should create new eval cases; it should not directly rewrite production behavior without tests.",
                    ],
                ),
            ],
        },
    ]

    for idx, spec in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        add_title(slide, spec["title"], spec.get("subtitle"), accent=spec.get("accent", BLUE))

        if spec.get("type") == "cover":
            slide.shapes.add_picture(str(TREE_IMAGE), Inches(0.82), Inches(1.38), height=Inches(5.42))
            add_section_label(slide, 10.1, 1.43, "7-layer map", BLUE)
            add_bullets(
                slide,
                10.02,
                1.92,
                2.58,
                3.8,
                [
                    "Data and catalog facts are deterministic.",
                    "Retrieval finds candidates before generation.",
                    "Decision logic chooses what may be said.",
                    "Prompts only write or reason inside bounded evidence.",
                    "Verification and feedback close the loop.",
                ],
                font_size=11.5,
            )
        else:
            add_tree_thumb(slide)
            cards = spec["cards"]
            add_card(slide, 0.58, 1.45, 3.75, 5.36, cards[0][0], cards[0][1], accent=spec["accent"])
            add_card(slide, 4.48, 1.45, 3.55, 2.55, cards[1][0], cards[1][1], accent=spec["accent"])
            add_card(slide, 4.48, 4.26, 3.55, 2.55, cards[2][0], cards[2][1], accent=spec["accent"])

        add_footer(slide, idx)

    return prs


def main() -> None:
    if not TREE_IMAGE.exists():
        raise FileNotFoundError(f"Missing required image: {TREE_IMAGE}")
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
