from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE_TREE = ROOT / "docs" / "assets" / "inventory_pipeline_tree.png"
CROP_DIR = ROOT / "docs" / "assets" / "rag_pipeline_slide_crops"
OUTPUT = ROOT / "docs" / "inventory_rag_pipeline_engineering_7_slides.pptx"


INK = RGBColor(20, 28, 44)
MUTED = RGBColor(84, 96, 116)
LINE = RGBColor(220, 226, 235)
BLUE = RGBColor(32, 95, 166)
GREEN = RGBColor(31, 132, 82)
ORANGE = RGBColor(205, 103, 17)
PURPLE = RGBColor(105, 102, 198)
MAGENTA = RGBColor(170, 81, 142)
GRAY = RGBColor(92, 104, 124)
WHITE = RGBColor(255, 255, 255)


SLIDES = [
    {
        "title": "RAG Pipeline Scope",
        "subtitle": "Only the retrieval-augmented answering path: question -> evidence -> grounded answer.",
        "accent": BLUE,
        "crop": "runtime_flow",
        "layout": "wide",
        "sections": [
            (
                "What It Does",
                [
                    "Converts a customer question into searchable constraints.",
                    "Retrieves catalog evidence before writing any answer.",
                    "Ranks evidence, builds an answer contract, then verifies the final reply.",
                ],
            ),
            (
                "Technology",
                [
                    "Python orchestration inside InventoryService.",
                    "Structured catalog records, vector search, lexical search, custom ranking.",
                    "Optional LLM prompts only after evidence is bounded.",
                ],
            ),
            (
                "Engineering Logic",
                [
                    "Hard product facts come from catalog fields, not model memory.",
                    "The pipeline fails closed: clarify or abstain when evidence is weak.",
                    "Feedback creates evaluation cases; it should not directly mutate answers.",
                ],
            ),
        ],
    },
    {
        "title": "DB Setup + Catalog Ground Truth",
        "subtitle": "The catalog is the source of truth for product facts used by RAG.",
        "accent": GREEN,
        "crop": "catalog_db",
        "sections": [
            (
                "What It Does",
                [
                    "Stores every product as structured facts: product_id, SKU, name, category, price, stock, status.",
                    "Keeps searchable attributes: color, size, fabric, gender, occasion, design_id, brand, tags.",
                    "Defines what the bot is allowed to know about inventory.",
                ],
            ),
            (
                "Technology",
                [
                    "JSONL catalog: data/inventory/catalog.jsonl.",
                    "Pydantic models validate records before use.",
                    "Optional SQLite mirror for operational lookup and sync workflows.",
                ],
            ),
            (
                "Engineering Logic",
                [
                    "No LLM is used at this layer.",
                    "include_in_rag and status control whether a product can appear in RAG answers.",
                    "Exact fields beat prose: stock, price, SKU, size, and product id are never guessed.",
                ],
            ),
        ],
    },
    {
        "title": "Query Understanding",
        "subtitle": "Turns Bangla, Banglish, or English text into intent and shopping slots.",
        "accent": PURPLE,
        "crop": "understanding",
        "sections": [
            (
                "What It Does",
                [
                    "Normalizes customer text and expands common Banglish aliases.",
                    "Classifies intent: search, compare, size availability, same-design color, styling, policy, or order.",
                    "Extracts slots: category, color, size, fabric, budget, occasion, gender, brand, in-stock requirement.",
                ],
            ),
            (
                "Technology",
                [
                    "Regex patterns, alias dictionaries, Bangla digit normalization, fuzzy correction.",
                    "Optional Ollama qwen3:8b JSON prompts for classifier, slot extractor, and conversation planner.",
                    "Recent conversation state helps resolve follow-ups like 'same design blue e ache?'.",
                ],
            ),
            (
                "Engineering Logic",
                [
                    "Deterministic extraction wins for hard fields.",
                    "LLM output fills gaps but should not override safer exact matches.",
                    "Low-confidence understanding triggers one clarification question.",
                ],
            ),
        ],
    },
    {
        "title": "Structured Retrieval",
        "subtitle": "Exact filters answer hard inventory questions before semantic search gets involved.",
        "accent": ORANGE,
        "crop": "structured_retrieval",
        "sections": [
            (
                "What It Does",
                [
                    "Filters catalog records by category, color, size, price, stock, gender, fabric, occasion, design_id.",
                    "Handles exact product questions like size availability and same-design different-color checks.",
                    "Returns precise candidates when the user gives concrete constraints.",
                ],
            ),
            (
                "Technology",
                [
                    "Python filtering over validated Pydantic inventory records.",
                    "Retail-specific logic in fashion_retail.py and inventory_service.py.",
                    "Structured metadata from catalog attributes drives filtering.",
                ],
            ),
            (
                "Engineering Logic",
                [
                    "Use structured retrieval for correctness-critical fields.",
                    "design_id/variant groups identify same design across colors or sizes.",
                    "Out-of-stock products may be mentioned only with a clear caveat.",
                ],
            ),
        ],
    },
    {
        "title": "Vector + Lexical Retrieval",
        "subtitle": "Semantic and keyword search increase recall when the customer wording is fuzzy.",
        "accent": ORANGE,
        "crop": "vector_lexical",
        "sections": [
            (
                "What It Does",
                [
                    "Embeds product search text and the customer query for semantic matching.",
                    "Uses lexical matching for product names, SKU, brand, category, and exact keywords.",
                    "Combines semantic recall with metadata filters so fuzzy wording still lands on real products.",
                ],
            ),
            (
                "Technology",
                [
                    "Embedding providers: transformers/multilingual, OpenAI-compatible, or deterministic fallback.",
                    "Vector stores: local JSONL by default; Elasticsearch adapter supports dense vectors and filters.",
                    "BM25 or Elasticsearch multi_match supports lexical retrieval.",
                ],
            ),
            (
                "Engineering Logic",
                [
                    "Embeddings help recall; they do not override exact catalog fields.",
                    "Lexical search protects exact names and SKUs that embeddings may blur.",
                    "Retrieved hits are candidates only; ranking and evidence gates still decide final claims.",
                ],
            ),
        ],
    },
    {
        "title": "Ranking + Evidence Contract",
        "subtitle": "Transforms candidates into allowed facts and controlled answer decisions.",
        "accent": MAGENTA,
        "crop": "decisions",
        "sections": [
            (
                "What It Does",
                [
                    "Scores candidates by relevance, exact slot match, stock fit, price fit, category fit, and product-name match.",
                    "Builds an evidence contract: primary item, alternatives, cross-sells, exclusions, tradeoffs, missing facts.",
                    "Decides whether to answer, clarify, or abstain.",
                ],
            ),
            (
                "Technology",
                [
                    "Custom ecommerce reranker and deterministic decisioning code.",
                    "Pydantic evidence contract objects.",
                    "Optional bounded LLM reasoner can choose only from retrieved product ids.",
                ],
            ),
            (
                "Engineering Logic",
                [
                    "Only supported catalog facts can pass into the final answer.",
                    "Cross-sells are add-ons, not substitutes.",
                    "Weak evidence, missing required facts, or contradictions trigger clarification/abstention.",
                ],
            ),
        ],
    },
    {
        "title": "Answer Writing + Verification",
        "subtitle": "Writes the customer reply from evidence, then checks it before returning.",
        "accent": GRAY,
        "crop": "answer_verify",
        "sections": [
            (
                "What It Does",
                [
                    "Writes a concise answer in the customer's language style: Bangla, Banglish, or English.",
                    "Mentions only approved products, prices, stock, caveats, and follow-up questions.",
                    "Verifies that the final text did not invent product facts or override the answer plan.",
                ],
            ),
            (
                "Technology",
                [
                    "Deterministic templates for safe fallback answers.",
                    "Optional Ollama/OpenAI-compatible natural writer with a strict ecommerce writer prompt.",
                    "Optional answer critic checks the reply against product facts.",
                ],
            ),
            (
                "Engineering Logic",
                [
                    "The writer expresses the decision; it does not choose products.",
                    "writer_contract and answer_plan control allowed claims.",
                    "If verification fails, the system falls back or abstains instead of shipping a fluent wrong answer.",
                ],
            ),
        ],
    },
]


CROPS = {
    "runtime_flow": (70, 1260, 2925, 1505),
    "catalog_db": (60, 470, 535, 915),
    "understanding": (1015, 470, 1465, 1210),
    "structured_retrieval": (1490, 470, 1945, 920),
    "vector_lexical": (1490, 905, 1945, 1210),
    "decisions": (1980, 470, 2410, 1210),
    "answer_verify": (2460, 470, 2905, 1210),
}


def _clean_background(image: Image.Image) -> Image.Image:
    image = image.convert("RGB")
    pixels = image.load()
    width, height = image.size
    for y in range(height):
        for x in range(width):
            r, g, b = pixels[x, y]
            if r > 232 and g > 235 and b > 240:
                pixels[x, y] = (255, 255, 255)
    return image


def make_crops() -> dict[str, Path]:
    if not SOURCE_TREE.exists():
        raise FileNotFoundError(f"Missing source image: {SOURCE_TREE}")
    CROP_DIR.mkdir(parents=True, exist_ok=True)
    source = Image.open(SOURCE_TREE).convert("RGB")
    output_paths: dict[str, Path] = {}
    for name, box in CROPS.items():
        crop = source.crop(box)
        crop = _clean_background(crop)
        draw = ImageDraw.Draw(crop)
        draw.rounded_rectangle(
            (3, 3, crop.width - 4, crop.height - 4),
            radius=18,
            outline=(218, 226, 238),
            width=3,
        )
        path = CROP_DIR / f"{name}.png"
        crop.save(path)
        output_paths[name] = path
    return output_paths


def add_textbox(slide, x, y, w, h, text="", *, size=16, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str, accent: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_textbox(slide, 0.55, 0.28, 8.7, 0.38, title, size=24, bold=True)
    add_textbox(slide, 0.58, 0.74, 10.9, 0.32, subtitle, size=10.5, color=MUTED)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.11), Inches(12.18), Inches(0.035))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent


def add_bullets(slide, x, y, w, h, bullets: list[str], *, size=10.8) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(5.5)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p._p.get_or_add_pPr().set("marL", "210000")
        p._p.get_or_add_pPr().set("indent", "-130000")


def add_section(slide, x, y, w, h, title: str, bullets: list[str], accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.1)
    add_textbox(slide, x + 0.17, y + 0.12, w - 0.34, 0.26, title, size=12.2, bold=True, color=accent)
    add_bullets(slide, x + 0.17, y + 0.43, w - 0.34, h - 0.5, bullets)


def add_footer(slide, slide_no: int) -> None:
    add_textbox(
        slide,
        0.58,
        7.12,
        7.2,
        0.2,
        "RAG pipeline only: catalog facts -> retrieval -> evidence -> verified answer",
        size=8.2,
        color=MUTED,
    )
    add_textbox(slide, 12.1, 7.12, 0.7, 0.2, str(slide_no), size=8.2, color=MUTED, align=PP_ALIGN.RIGHT)


def add_image_with_fit(slide, image_path: Path, x: float, y: float, max_w: float, max_h: float) -> None:
    with Image.open(image_path) as img:
        width, height = img.size
    ratio = width / height
    target_w = max_w
    target_h = target_w / ratio
    if target_h > max_h:
        target_h = max_h
        target_w = target_h * ratio
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(target_w), height=Inches(target_h))


def build_deck(crop_paths: dict[str, Path]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, spec in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        add_title(slide, spec["title"], spec["subtitle"], spec["accent"])

        if spec.get("layout") == "wide":
            add_image_with_fit(slide, crop_paths[spec["crop"]], 0.62, 1.45, 12.1, 1.65)
            x_positions = [0.62, 4.55, 8.48]
            for x, section in zip(x_positions, spec["sections"], strict=True):
                add_section(slide, x, 3.42, 3.65, 3.25, section[0], section[1], spec["accent"])
        else:
            add_image_with_fit(slide, crop_paths[spec["crop"]], 0.65, 1.42, 4.55, 5.35)
            add_section(slide, 5.55, 1.42, 3.35, 2.48, spec["sections"][0][0], spec["sections"][0][1], spec["accent"])
            add_section(slide, 9.18, 1.42, 3.35, 2.48, spec["sections"][1][0], spec["sections"][1][1], spec["accent"])
            add_section(slide, 5.55, 4.24, 6.98, 2.53, spec["sections"][2][0], spec["sections"][2][1], spec["accent"])

        add_footer(slide, idx)

    return prs


def main() -> None:
    crop_paths = make_crops()
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck(crop_paths)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
